# Copyright(C) 2025-2026 Advanced Micro Devices, Inc. All rights reserved.
# SPDX-License-Identifier: MIT
"""
End-to-end integration test for PPTX extraction through the full RAG + LLM
pipeline.

Creates a test PPTX programmatically (no sensitive files), indexes it via
RAGSDK, queries via rag.query() (real LLM on Lemonade), and verifies the
answers contain expected content.

Skips when Lemonade is not reachable.
"""

from __future__ import annotations

import time
from pathlib import Path
from urllib.error import URLError
from urllib.request import urlopen

import pytest


def _lemonade_available(base_url="http://localhost:13305"):
    """Check if Lemonade server is reachable."""
    try:
        with urlopen(f"{base_url}/api/v1/models", timeout=3) as r:
            return r.status == 200, "ok"
    except (URLError, OSError) as e:
        return False, f"Lemonade not reachable: {e}"


def _create_test_pptx(path: Path) -> None:
    """Build a multi-slide PPTX with known content for Q&A verification."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Aurora: Next-Gen Solar Panel Technology"
    slide1.placeholders[1].text = "Q3 2026 Technical Review"

    # Slide 2: Problem statement
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Current Challenges"
    slide2.placeholders[1].text = (
        "Traditional silicon solar panels have reached a theoretical efficiency "
        "ceiling of 29%. Manufacturing costs remain high at $0.30 per watt. "
        "Degradation rates of 0.5% per year reduce lifetime output by 12% over "
        "25 years. Project Aurora addresses these limitations with perovskite "
        "tandem cell technology."
    )

    # Slide 3: Solution with specific data
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Aurora Perovskite Solution"
    slide3.placeholders[1].text = (
        "The Aurora tandem cell combines perovskite and silicon layers to achieve "
        "33.7% efficiency in lab conditions. Manufacturing cost target is $0.18 "
        "per watt using roll-to-roll printing. The perovskite layer absorbs "
        "blue and green wavelengths while silicon handles red and infrared. "
        "Field trials in Phoenix, Arizona showed 31.2% real-world efficiency."
    )

    # Slide 4: Timeline with speaker notes
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Project Timeline"
    slide4.placeholders[1].text = (
        "Phase 1 (Q1 2026): Lab prototype validation. "
        "Phase 2 (Q3 2026): Pilot manufacturing line in Dresden, Germany. "
        "Phase 3 (Q1 2027): Commercial production at 500MW annual capacity. "
        "Total investment: $47 million across all phases."
    )
    slide4.notes_slide.notes_text_frame.text = (
        "The Dresden facility was chosen for its existing semiconductor "
        "supply chain and skilled workforce. EU Green Deal subsidies cover "
        "30% of Phase 2 costs."
    )

    # Slide 5: Team with a table
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    tx = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.5))
    tx.text_frame.text = "Core Team"
    table_shape = slide5.shapes.add_table(
        4, 3, Inches(0.5), Inches(1), Inches(8), Inches(3)
    )
    table = table_shape.table
    for c, text in enumerate(["Name", "Role", "Location"]):
        table.cell(0, c).text = text
    for r, row in enumerate(
        [
            ["Dr. Elena Vasquez", "Chief Scientist", "MIT"],
            ["Marcus Chen", "Manufacturing Lead", "Dresden"],
            ["Aisha Patel", "Field Testing Director", "Phoenix"],
        ],
        1,
    ):
        for c, text in enumerate(row):
            table.cell(r, c).text = text

    prs.save(str(path))


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


@pytest.mark.integration
class TestPptxRagE2E:
    """End-to-end PPTX → RAG → LLM tests against real Lemonade."""

    @pytest.fixture(autouse=True)
    def _skip_without_lemonade(self):
        ok, reason = _lemonade_available()
        if not ok:
            pytest.skip(f"Skipping E2E: {reason}")

    @pytest.fixture
    def pptx_path(self, tmp_path):
        path = tmp_path / "aurora_solar.pptx"
        _create_test_pptx(path)
        return str(path)

    @pytest.fixture
    def rag(self, tmp_path, pptx_path):
        from gaia.rag.sdk import RAGSDK, RAGConfig

        config = RAGConfig(
            cache_dir=str(tmp_path / ".rag_cache"),
            show_stats=False,
            vlm_model="Gemma-4-E4B-it-GGUF",
        )
        rag = RAGSDK(config=config)
        stats = rag.index_document(pptx_path)
        assert stats["success"], f"Indexing failed: {stats.get('error')}"
        assert stats["num_chunks"] > 0
        return rag

    def test_index_creates_chunks(self, rag):
        """Indexing a PPTX creates a non-empty vector index."""
        assert len(rag.chunks) >= 1

    def test_retrieval_finds_efficiency_data(self, rag):
        """Semantic search retrieves the correct chunk for efficiency questions."""
        chunks, scores = rag._retrieve_chunks(
            "What efficiency does the Aurora panel achieve?"
        )
        assert chunks, "No chunks retrieved"
        top = chunks[0].lower()
        assert "33.7" in top or "efficiency" in top

    def test_llm_answers_cost_question(self, rag):
        """Full LLM Q&A returns the manufacturing cost from the slides."""
        start = time.time()
        response = rag.query("What is the manufacturing cost target for Aurora?")
        latency = time.time() - start

        answer = response.text if hasattr(response, "text") else str(response)
        answer_lower = answer.lower()

        # The slides say "$0.18 per watt"
        assert (
            "$0.18" in answer_lower or "0.18" in answer_lower
        ), f"Expected '$0.18' in answer but got: {answer[:300]}"
        assert latency < 120, f"LLM query took too long: {latency:.1f}s"

    def test_llm_answers_location_question(self, rag):
        """LLM correctly identifies the manufacturing location."""
        response = rag.query("Where is the pilot manufacturing facility located?")
        answer = response.text if hasattr(response, "text") else str(response)

        assert (
            "dresden" in answer.lower()
        ), f"Expected 'Dresden' in answer but got: {answer[:300]}"

    def test_llm_answers_team_question(self, rag):
        """LLM retrieves team member info from the table."""
        response = rag.query("Who is the Chief Scientist on the project?")
        answer = response.text if hasattr(response, "text") else str(response)

        assert (
            "vasquez" in answer.lower() or "elena" in answer.lower()
        ), f"Expected 'Vasquez' or 'Elena' in answer but got: {answer[:300]}"

    def test_speaker_notes_indexed(self, rag):
        """Speaker notes content is retrievable."""
        chunks, _ = rag._retrieve_chunks("EU Green Deal subsidies")
        assert chunks, "No chunks retrieved for speaker notes query"
        found = any("green deal" in c.lower() or "subsid" in c.lower() for c in chunks)
        assert found, "Speaker notes content not found in any chunk"
