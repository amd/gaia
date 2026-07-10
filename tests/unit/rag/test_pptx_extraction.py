# Copyright(C) 2025-2026 Advanced Micro Devices, Inc. All rights reserved.
# SPDX-License-Identifier: MIT
"""
Unit tests for PowerPoint (.pptx) extraction in gaia.rag.sdk.

Covers:
- Text extraction from shapes, titles, tables, speaker notes
- Multi-slide [Page N] marker generation
- Image extraction with VLM (mocked)
- Graceful handling of empty / corrupted presentations
- Metadata structure

Fixtures are built programmatically with python-pptx so the tests remain
hermetic and don't require committing binary fixture files to the repo.
"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

pptx = pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

from gaia.rag.sdk import RAGSDK, RAGConfig  # noqa: E402

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def rag(tmp_path: Path) -> RAGSDK:
    """
    A RAGSDK instance scoped to tmp_path with heavy ML deps stubbed out.

    Same pattern as test_pdf_extraction_errors.py — stubs sentence-transformers,
    faiss, VLM, and chat/LLM initialization.
    """
    config = RAGConfig(
        cache_dir=str(tmp_path / ".gaia"),
        show_stats=False,
        use_local_llm=False,
    )

    fake_vlm = MagicMock(name="VLMClient")
    fake_vlm.check_availability.return_value = False

    with (
        patch.object(RAGSDK, "_check_dependencies", return_value=None),
        patch("gaia.rag.sdk.AgentSDK", autospec=True) as mock_agent_sdk,
        patch("gaia.llm.VLMClient", return_value=fake_vlm),
    ):
        mock_agent_sdk.return_value = MagicMock(name="AgentSDK")
        instance = RAGSDK(config=config)

    instance._test_vlm_patch = patch("gaia.llm.VLMClient", return_value=fake_vlm)
    instance._test_vlm_patch.start()
    yield instance
    instance._test_vlm_patch.stop()


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _create_pptx(path: Path, slides: list) -> None:
    """Create a test .pptx programmatically.

    Each dict in *slides* may have keys:

    - ``title`` (str): slide title text
    - ``body`` (str): body placeholder text
    - ``notes`` (str): speaker notes
    - ``table`` (list[list[str]]): rows for a table (first row is header)
    """
    prs = Presentation()
    for content in slides:
        slide_layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(slide_layout)

        if "title" in content and slide.shapes.title:
            slide.shapes.title.text = content["title"]

        if "body" in content:
            # Use the content placeholder (index 1)
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = content["body"]
                    break

        if "notes" in content:
            slide.notes_slide.notes_text_frame.text = content["notes"]

        if "table" in content:
            rows_data = content["table"]
            num_rows = len(rows_data)
            num_cols = len(rows_data[0]) if rows_data else 0
            if num_rows and num_cols:
                table_shape = slide.shapes.add_table(
                    num_rows, num_cols, Inches(1), Inches(3), Inches(6), Inches(2)
                )
                table = table_shape.table
                for r_idx, row in enumerate(rows_data):
                    for c_idx, cell_text in enumerate(row):
                        table.cell(r_idx, c_idx).text = cell_text

    prs.save(str(path))


def _create_pptx_with_image(path: Path, image_bytes: bytes) -> None:
    """Create a .pptx with a single slide containing an embedded image."""
    import io

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(1), Inches(1))
    prs.save(str(path))


def _make_red_png() -> bytes:
    """Generate a minimal 10x10 red PNG image."""
    from PIL import Image

    img = Image.new("RGB", (10, 10), color="red")
    import io

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Tests — Text Extraction
# ---------------------------------------------------------------------------


class TestPptxTextExtraction:
    """Tests for _extract_text_from_pptx text handling."""

    def test_basic_text_extraction(self, rag, tmp_path):
        """Single slide with title + body — both appear in output."""
        pptx_path = tmp_path / "basic.pptx"
        _create_pptx(pptx_path, [{"title": "Hello World", "body": "Test content here"}])

        text, num_slides, metadata = rag._extract_text_from_pptx(str(pptx_path))

        assert "Hello World" in text
        assert "Test content here" in text
        assert num_slides == 1

    def test_multi_slide_page_markers(self, rag, tmp_path):
        """Multiple slides produce [Page 1], [Page 2], [Page 3] markers."""
        pptx_path = tmp_path / "multi.pptx"
        _create_pptx(
            pptx_path,
            [
                {"title": "Slide One", "body": "Content 1"},
                {"title": "Slide Two", "body": "Content 2"},
                {"title": "Slide Three", "body": "Content 3"},
            ],
        )

        text, num_slides, _ = rag._extract_text_from_pptx(str(pptx_path))

        assert num_slides == 3
        assert "[Page 1]" in text
        assert "[Page 2]" in text
        assert "[Page 3]" in text
        assert "Slide One" in text
        assert "Slide Two" in text
        assert "Slide Three" in text

    def test_speaker_notes_extraction(self, rag, tmp_path):
        """Speaker notes appear in extracted text."""
        pptx_path = tmp_path / "notes.pptx"
        _create_pptx(
            pptx_path,
            [{"title": "Presentation", "notes": "Remember to mention the deadline"}],
        )

        text, _, _ = rag._extract_text_from_pptx(str(pptx_path))

        assert "Speaker Notes:" in text
        assert "Remember to mention the deadline" in text

    def test_table_extraction(self, rag, tmp_path):
        """Tables are extracted as markdown by the pure-python (python-pptx) path.

        The PPTX→PDF fast path (LibreOffice/PowerPoint COM) flattens tables to
        plain text — markdown table structure is produced only by the pure-python
        fallback. Force that fallback here so the assertion is deterministic
        regardless of whether LibreOffice is installed on the runner.
        """
        pptx_path = tmp_path / "table.pptx"
        _create_pptx(
            pptx_path,
            [
                {
                    "title": "Data Table",
                    "table": [
                        ["Name", "Age", "City"],
                        ["Alice", "30", "Boston"],
                        ["Bob", "25", "Denver"],
                    ],
                }
            ],
        )

        # Disable PPTX→PDF conversion so the pure-python markdown extractor runs.
        with patch(
            "gaia.rag.pptx_utils.convert_pptx_to_pdf",
            side_effect=RuntimeError("PDF conversion disabled for this test"),
        ):
            text, _, _ = rag._extract_text_from_pptx(str(pptx_path))

        assert "Alice" in text
        assert "Bob" in text
        assert "Boston" in text
        assert "Denver" in text
        # Should have markdown table formatting
        assert "|" in text
        assert "---" in text


# ---------------------------------------------------------------------------
# Tests — Edge Cases
# ---------------------------------------------------------------------------


class TestPptxEdgeCases:
    """Tests for empty, blank, and corrupted presentations."""

    def test_empty_presentation_raises(self, rag, tmp_path):
        """PPTX with no slides raises ValueError with guidance."""
        pptx_path = tmp_path / "empty.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        # Empty presentation has 0 slides — should not crash
        # but the result may have no content. The method checks
        # has_any_content after the loop.
        # With 0 slides, pages_data is empty, has_any_content is False
        # unless total_slides==0 means any() on empty list returns False
        with pytest.raises(ValueError, match="No extractable text"):
            rag._extract_text_from_pptx(str(pptx_path))

    def test_blank_slide(self, rag, tmp_path):
        """PPTX with one completely blank slide raises ValueError."""
        pptx_path = tmp_path / "blank.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        prs.save(str(pptx_path))

        with pytest.raises(ValueError, match="No extractable text"):
            rag._extract_text_from_pptx(str(pptx_path))

    def test_corrupted_file_raises(self, rag, tmp_path):
        """Garbage bytes with .pptx extension raises ValueError."""
        pptx_path = tmp_path / "corrupt.pptx"
        pptx_path.write_bytes(b"this is not a pptx file at all")

        with pytest.raises(ValueError, match="corrupted|not a valid"):
            rag._extract_text_from_pptx(str(pptx_path))


# ---------------------------------------------------------------------------
# Tests — Metadata
# ---------------------------------------------------------------------------


class TestPptxMetadata:
    """Tests for metadata structure from _extract_text_from_pptx."""

    def test_metadata_keys(self, rag, tmp_path):
        """Metadata contains expected keys and values."""
        pptx_path = tmp_path / "meta.pptx"
        _create_pptx(pptx_path, [{"title": "Test", "body": "Content"}])

        _, num_slides, metadata = rag._extract_text_from_pptx(str(pptx_path))

        assert num_slides == 1
        assert metadata["num_slides"] == 1
        assert metadata["vlm_slides"] == 0
        assert metadata["total_images"] == 0
        assert metadata["vlm_checked"] is True
        assert metadata["vlm_available"] is False
        assert metadata["pptx_status"] == "readable"


# ---------------------------------------------------------------------------
# Tests — VLM Integration (mocked)
# ---------------------------------------------------------------------------


class TestPptxVLMIntegration:
    """Tests for VLM image extraction from slides."""

    def test_vlm_not_available_skips_images(self, rag, tmp_path):
        """When VLM is unavailable, embedded images are not processed.

        We add text alongside the image so the empty-content check doesn't
        fire — the test is specifically about VLM skipping, not about empty
        presentations.
        """
        PIL = pytest.importorskip("PIL")  # noqa: N806, F841
        import io as _io

        pptx_path = tmp_path / "with_image.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(_io.BytesIO(_make_red_png()), Inches(1), Inches(1))
        # Add a text box so the slide is not considered empty
        tx_box = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(3), Inches(1))
        tx_box.text_frame.text = "Slide with image and text"
        prs.save(str(pptx_path))

        text, _, metadata = rag._extract_text_from_pptx(str(pptx_path))

        # VLM is mocked as unavailable in the fixture
        assert metadata["total_images"] == 0
        assert metadata["vlm_slides"] == 0
        assert "Slide with image and text" in text

    def test_vlm_processes_images(self, rag, tmp_path):
        """When VLM is available, extracted image text appears in output."""
        PIL = pytest.importorskip("PIL")  # noqa: N806, F841
        import io as _io

        pptx_path = tmp_path / "vlm_test.pptx"
        # Create slide with both an image and text (so empty-check doesn't fire)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(_io.BytesIO(_make_red_png()), Inches(1), Inches(1))
        tx_box = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(3), Inches(1))
        tx_box.text_frame.text = "Slide with VLM image"
        prs.save(str(pptx_path))

        # Patch VLM to be available and return text
        fake_vlm = MagicMock(name="VLMClient")
        fake_vlm.check_availability.return_value = True
        fake_vlm.extract_from_page_images.return_value = [
            {
                "image_num": 1,
                "text": "Extracted text from slide image",
                "dimensions": "10x10",
                "size_kb": 0.5,
            }
        ]

        # Need to re-patch VLM for this specific test
        rag._test_vlm_patch.stop()
        try:
            with patch("gaia.llm.VLMClient", return_value=fake_vlm):
                text, _, metadata = rag._extract_text_from_pptx(str(pptx_path))

            assert "Extracted text from slide image" in text
            assert metadata["total_images"] == 1
            assert metadata["vlm_slides"] == 1
        finally:
            # Re-start the original patch for fixture cleanup
            rag._test_vlm_patch.start()


# ---------------------------------------------------------------------------
# Tests — Integration with _extract_text_from_file dispatcher
# ---------------------------------------------------------------------------


class TestPptxDispatcher:
    """Tests for .pptx routing through _extract_text_from_file."""

    def test_pptx_dispatches_correctly(self, rag, tmp_path):
        """_extract_text_from_file routes .pptx to _extract_text_from_pptx."""
        pptx_path = tmp_path / "dispatch.pptx"
        _create_pptx(pptx_path, [{"title": "Dispatch Test", "body": "Works"}])

        text, metadata = rag._extract_text_from_file(str(pptx_path))

        assert "Dispatch Test" in text
        assert "Works" in text
        assert metadata["num_pages"] == 1  # num_slides mapped to num_pages


# ---------------------------------------------------------------------------
# Tests — pptx_utils module
# ---------------------------------------------------------------------------


class TestPptxUtils:
    """Direct tests for pptx_utils helper functions."""

    def test_extract_text_from_slide(self):
        """extract_text_from_slide captures text frame content."""
        from gaia.rag.pptx_utils import extract_text_from_slide

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "My Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = "Body text"
                break

        text = extract_text_from_slide(slide, slide_num=1)
        assert "My Title" in text
        assert "Body text" in text

    def test_extract_notes_from_slide(self):
        """extract_notes_from_slide returns speaker notes."""
        from gaia.rag.pptx_utils import extract_notes_from_slide

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.notes_slide.notes_text_frame.text = "These are my notes"

        notes = extract_notes_from_slide(slide)
        assert notes == "These are my notes"

    def test_extract_notes_empty(self):
        """extract_notes_from_slide returns empty string when no notes."""
        from gaia.rag.pptx_utils import extract_notes_from_slide

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        notes = extract_notes_from_slide(slide)
        assert notes == ""

    def test_count_images_in_slide_no_images(self):
        """count_images_in_slide returns (False, 0) for text-only slides."""
        from gaia.rag.pptx_utils import count_images_in_slide

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "No images here"

        has_images, count = count_images_in_slide(slide)
        assert has_images is False
        assert count == 0

    def test_count_images_in_slide_with_image(self):
        """count_images_in_slide detects embedded images."""
        PIL = pytest.importorskip("PIL")  # noqa: N806, F841
        import io

        from gaia.rag.pptx_utils import count_images_in_slide

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(io.BytesIO(_make_red_png()), Inches(1), Inches(1))

        has_images, count = count_images_in_slide(slide)
        assert has_images is True
        assert count >= 1

    def test_table_to_markdown(self):
        """_table_to_markdown produces valid markdown table."""
        from gaia.rag.pptx_utils import _table_to_markdown

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table_shape = slide.shapes.add_table(
            2, 3, Inches(1), Inches(1), Inches(6), Inches(2)
        )
        table = table_shape.table
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Age"
        table.cell(0, 2).text = "City"
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "30"
        table.cell(1, 2).text = "Boston"

        md = _table_to_markdown(table)
        assert "| Name | Age | City |" in md
        assert "| --- | --- | --- |" in md
        assert "| Alice | 30 | Boston |" in md


# ---------------------------------------------------------------------------
# Tests — PowerPoint COM → PDF conversion path
# ---------------------------------------------------------------------------


class TestPptxPdfConversion:
    """Tests for the PPTX → PDF → existing pipeline fast path."""

    def test_pdf_conversion_used_when_available(self, rag, tmp_path):
        """When convert_pptx_to_pdf returns a PDF, _extract_text_from_pdf is used."""
        pptx_path = tmp_path / "conv.pptx"
        _create_pptx(
            pptx_path,
            [
                {"title": "Converted Slide", "body": "Body text", "notes": "My notes"},
            ],
        )

        fake_pdf_text = "[Page 1]\nConverted slide content from PDF"
        fake_pdf_metadata = {
            "num_pages": 1,
            "vlm_pages": 0,
            "total_images": 0,
            "vlm_checked": True,
            "vlm_available": False,
            "pdf_status": "readable",
        }

        with (
            patch(
                "gaia.rag.pptx_utils.convert_pptx_to_pdf",
                return_value="/fake/output.pdf",
            ),
            patch.object(
                rag,
                "_extract_text_from_pdf",
                return_value=(fake_pdf_text, 1, fake_pdf_metadata),
            ) as mock_pdf,
        ):
            text, num_slides, metadata = rag._extract_text_from_pptx(str(pptx_path))

        mock_pdf.assert_called_once_with("/fake/output.pdf")
        assert "Converted slide content from PDF" in text
        assert "My notes" in text  # Notes appended from python-pptx
        assert metadata["conversion"] == "powerpoint_com"

    def test_pdf_conversion_fallback_on_none(self, rag, tmp_path):
        """When convert_pptx_to_pdf returns None, python-pptx fallback runs."""
        pptx_path = tmp_path / "fallback.pptx"
        _create_pptx(
            pptx_path, [{"title": "Fallback Slide", "body": "Fallback content"}]
        )

        with patch("gaia.rag.pptx_utils.convert_pptx_to_pdf", return_value=None):
            text, num_slides, metadata = rag._extract_text_from_pptx(str(pptx_path))

        assert "Fallback Slide" in text
        assert "Fallback content" in text
        assert "conversion" not in metadata  # No conversion key in fallback path

    def test_pdf_conversion_fallback_on_exception(self, rag, tmp_path):
        """When PDF extraction from converted file fails, python-pptx fallback runs."""
        pptx_path = tmp_path / "exc.pptx"
        _create_pptx(pptx_path, [{"title": "Exception Slide", "body": "Safe content"}])

        with (
            patch(
                "gaia.rag.pptx_utils.convert_pptx_to_pdf",
                return_value="/fake/output.pdf",
            ),
            patch.object(
                rag,
                "_extract_text_from_pdf",
                side_effect=RuntimeError("PDF parsing exploded"),
            ),
        ):
            text, num_slides, metadata = rag._extract_text_from_pptx(str(pptx_path))

        # Should have fallen back to python-pptx
        assert "Exception Slide" in text
        assert "Safe content" in text
        assert "conversion" not in metadata
