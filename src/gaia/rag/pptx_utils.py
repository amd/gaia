#!/usr/bin/env python3
# Copyright(C) 2025-2026 Advanced Micro Devices, Inc. All rights reserved.
# SPDX-License-Identifier: MIT

"""
PowerPoint (.pptx) extraction utilities for multi-modal RAG.

Extracts text, tables, speaker notes, and embedded images from PPTX slides.
Image processing mirrors ``pdf_utils.py`` (resize, compress, same dict format).
"""

from __future__ import annotations

import io
import logging
import os
import platform
import subprocess

logger = logging.getLogger(__name__)

# Shared constants (same as pdf_utils.py)
MAX_DIMENSION = 1600
MAX_SIZE_KB = 300
MAX_ITERATIONS = 5
MAX_GROUP_DEPTH = 5


def extract_images_from_slide(slide, slide_num: int) -> list[dict]:
    """
    Extract embedded images from a PPTX slide.

    Iterates picture shapes, extracts image blobs, and processes with PIL
    (resize / compress) identically to ``pdf_utils.extract_images_from_page_pymupdf``.

    Args:
        slide: A ``pptx.slide.Slide`` object.
        slide_num: Slide number (1-indexed, for logging).

    Returns:
        List of image dicts::

            [{"image_bytes": bytes, "width": int, "height": int,
              "format": "png", "size_kb": float}, ...]
    """
    images: list[dict] = []

    try:
        from PIL import Image
    except ImportError:
        logger.error("Pillow not installed. Install: uv pip install Pillow")
        return images

    for shape_index, shape in enumerate(_iter_shapes(slide.shapes)):
        if not hasattr(shape, "image"):
            continue

        try:
            image_blob = shape.image.blob

            img = Image.open(io.BytesIO(image_blob))

            width, height = img.size
            size_kb = len(image_blob) / 1024

            # Convert to RGB if needed
            if img.mode not in ("RGB", "RGBA"):
                logger.debug("Converting %s to RGB", img.mode)
                img = img.convert("RGB")

            # Resize if too large
            if width > MAX_DIMENSION or height > MAX_DIMENSION:
                scale = min(MAX_DIMENSION / width, MAX_DIMENSION / height)
                new_width = int(width * scale)
                new_height = int(height * scale)
                logger.info(
                    "   Resizing: %dx%d -> %dx%d", width, height, new_width, new_height
                )
                img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)

            # Save as optimized PNG
            png_buffer = io.BytesIO()
            img.save(png_buffer, format="PNG", optimize=True, compress_level=6)
            png_bytes = png_buffer.getvalue()
            size_kb = len(png_bytes) / 1024

            # Iteratively compress until target size is reached
            compression_iterations = 0
            while size_kb > MAX_SIZE_KB and compression_iterations < MAX_ITERATIONS:
                compression_iterations += 1
                logger.info(
                    "   Compressing (iteration %d): %.0fKB -> <%dKB",
                    compression_iterations,
                    size_kb,
                    MAX_SIZE_KB,
                )
                img = img.resize(
                    (img.width // 2, img.height // 2), Image.Resampling.LANCZOS
                )
                png_buffer = io.BytesIO()
                img.save(png_buffer, format="PNG", optimize=True, compress_level=9)
                png_bytes = png_buffer.getvalue()
                size_kb = len(png_bytes) / 1024

            if size_kb <= MAX_SIZE_KB:
                logger.info(
                    "   Compressed to %.0fKB (%dx%d) in %d iteration(s)",
                    size_kb,
                    img.width,
                    img.height,
                    compression_iterations,
                )
            else:
                logger.warning(
                    "   Could not compress below %dKB after %d iterations (final: %.0fKB)",
                    MAX_SIZE_KB,
                    MAX_ITERATIONS,
                    size_kb,
                )

            images.append(
                {
                    "image_bytes": png_bytes,
                    "width": img.width,
                    "height": img.height,
                    "format": "png",
                    "size_kb": size_kb,
                }
            )

            logger.debug(
                "Extracted image %d from slide %d: %dx%d, %.1fKB",
                shape_index + 1,
                slide_num,
                img.width,
                img.height,
                size_kb,
            )

        except Exception as e:
            # WMF/EMF metafiles, corrupt blobs, etc. — skip with warning.
            logger.warning(
                "Failed to extract image %d from slide %d: %s",
                shape_index + 1,
                slide_num,
                e,
            )
            continue

    return images


def count_images_in_slide(slide) -> tuple[bool, int]:
    """
    Fast check for embedded image presence without extraction.

    Args:
        slide: A ``pptx.slide.Slide`` object.

    Returns:
        ``(has_images, count)`` tuple.
    """
    count = 0
    for shape in _iter_shapes(slide.shapes):
        if hasattr(shape, "image"):
            count += 1
    return (count > 0, count)


def extract_text_from_slide(slide, slide_num: int) -> str:
    """
    Extract all native text from a PPTX slide.

    Handles text frames, tables (formatted as markdown), and group shapes
    (recursed up to :data:`MAX_GROUP_DEPTH` levels).

    Args:
        slide: A ``pptx.slide.Slide`` object.
        slide_num: Slide number (1-indexed, for logging).

    Returns:
        Concatenated slide text with paragraph boundaries.
    """
    logger.debug("Extracting text from slide %d", slide_num)
    parts: list[str] = []

    for shape in _iter_shapes(slide.shapes):
        # Table shapes — format as markdown
        if shape.has_table:
            table_md = _table_to_markdown(shape.table)
            if table_md:
                parts.append(table_md)
        # Text-bearing shapes (text boxes, titles, subtitles, etc.)
        elif shape.has_text_frame:
            text = _text_frame_to_str(shape.text_frame)
            if text:
                parts.append(text)

    return "\n\n".join(parts)


def extract_notes_from_slide(slide) -> str:
    """
    Extract speaker notes from a slide.

    Args:
        slide: A ``pptx.slide.Slide`` object.

    Returns:
        Notes text, or empty string if none.
    """
    try:
        if not slide.has_notes_slide:
            return ""
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            text = notes_slide.notes_text_frame.text
            return text.strip() if text else ""
    except Exception as e:
        logger.debug("Could not extract notes from slide: %s", e)
    return ""


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _iter_shapes(shapes, depth: int = 0):
    """Yield all leaf shapes, recursing into group shapes up to *MAX_GROUP_DEPTH*.

    Group shapes themselves are never yielded — only their children.  When the
    depth limit is reached the group is silently skipped (its children are not
    visited) to prevent stack overflow on pathological files.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # lazy import

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if depth < MAX_GROUP_DEPTH:
                yield from _iter_shapes(shape.shapes, depth + 1)
            else:
                logger.warning(
                    "Group shape nesting exceeds MAX_GROUP_DEPTH (%d); skipping children",
                    MAX_GROUP_DEPTH,
                )
        else:
            yield shape


def _text_frame_to_str(text_frame) -> str:
    """Join paragraphs from a text frame, preserving line breaks."""
    lines = [p.text.strip() for p in text_frame.paragraphs if p.text.strip()]
    return "\n".join(lines)


def _table_to_markdown(table) -> str:
    """Convert a ``pptx.table.Table`` to a markdown table string."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")

    if not rows:
        return ""

    # Insert header separator after first row
    num_cols = len(table.rows[0].cells)
    separator = "| " + " | ".join(["---"] * num_cols) + " |"
    rows.insert(1, separator)

    return "\n".join(rows)


# ---------------------------------------------------------------------------
# PPTX → PDF conversion (Windows only, requires PowerPoint)
# ---------------------------------------------------------------------------


def convert_pptx_to_pdf(pptx_path: str, output_dir: str) -> str | None:
    """Convert a PPTX file to PDF using PowerPoint COM automation.

    Only works on Windows with Microsoft PowerPoint installed.  Returns the
    path to the generated PDF on success, or ``None`` if conversion is not
    possible (wrong OS, PowerPoint not installed, timeout, etc.).

    This function never raises — the caller should fall back to python-pptx
    native extraction when ``None`` is returned.

    Args:
        pptx_path: Absolute path to the ``.pptx`` file.
        output_dir: Directory where the PDF will be written.

    Returns:
        Absolute path to the generated PDF, or ``None``.
    """
    if platform.system() != "Windows":
        logger.debug(
            "PPTX→PDF conversion requires Windows (current: %s)", platform.system()
        )
        return None

    from pathlib import Path  # already available, but keep import local for clarity

    pptx_abs = str(Path(pptx_path).resolve())
    pdf_name = Path(pptx_path).stem + ".pdf"
    pdf_abs = str(Path(output_dir).resolve() / pdf_name)

    # PowerShell script using PowerPoint COM.  Single-quoted paths handle
    # spaces correctly.  MsoTriState values are raw integers to avoid
    # needing the Office interop assembly.
    #   msoTrue = -1, msoFalse = 0, ppSaveAsPDF = 32
    ps_script = (
        "$ErrorActionPreference = 'Stop'; "
        "$ppt = New-Object -ComObject PowerPoint.Application; "
        "try { "
        f"  $pres = $ppt.Presentations.Open('{pptx_abs}', "
        "    [int]-1, "  # ReadOnly = msoTrue
        "    [int]0, "  # Untitled = msoFalse
        "    [int]0"  # WithWindow = msoFalse
        "  ); "
        f"  $pres.SaveAs('{pdf_abs}', 32); "  # 32 = ppSaveAsPDF
        "  $pres.Close(); "
        "} finally { "
        "  $ppt.Quit(); "
        "}"
    )

    try:
        result = subprocess.run(
            ["powershell", "-NoProfile", "-Command", ps_script],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )

        if result.returncode == 0 and os.path.exists(pdf_abs):
            logger.info("PPTX→PDF conversion succeeded: %s", pdf_abs)
            return pdf_abs

        logger.debug(
            "PPTX→PDF conversion failed (rc=%d): %s",
            result.returncode,
            result.stderr.strip()[:200] if result.stderr else "(no stderr)",
        )
        return None

    except subprocess.TimeoutExpired:
        logger.warning("PPTX→PDF conversion timed out after 120s")
        # Kill any orphaned PowerPoint process spawned by COM
        subprocess.run(
            ["taskkill", "/f", "/im", "POWERPNT.EXE"],
            capture_output=True,
            timeout=10,
            check=False,
        )
        return None
    except FileNotFoundError:
        logger.debug("PowerShell not found on PATH")
        return None
    except Exception as e:
        logger.debug("PPTX→PDF conversion failed: %s", e)
        return None
